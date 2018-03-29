#!/usr/bin/env/python3

import pandas as pd
import glob
import urllib

import requests
from pptx import Presentation
from pptx.util import Inches

known_textbooks = ['Knight, Physics for Scientists and Engineers: A Strategic Approach with Modern Physics, 3rd Edition',
                   'Knight, College Physics: A Strategic Approach, 2nd Edition']


def scientistparts(fname):
    """The scientist file looks like
    # Name
    Margaret Murnane
    # Description
    Words words words
    words words words

    # Sources
    sourcy source

    ...

    This is an iterator that returns the (item name, content)
    We expect to combine things (e.g. all of the lines in the 
    description) later. But not now, because, e.g., each line 
    in the textbook part is a separate textbook.
    """
    item_name,content = None,[]

    for line in open(fname):
        if line.startswith('#'):
            if item_name is not None:
                yield (item_name,[i for i in content if i.strip()])
            item_name = line.replace('#','').strip()
            content = []
        else:
            content.append(line.strip())
    if item_name is not None:
        yield(item_name,[i for i in content if i.strip()])

class Scientist:
    def __init__(self,fname):
        #self._scientistdict = dict(scientistparts(fname))
        # Can't quite do that because older format had you 
        # putting in two #Textbook parts instead of just two 
        # lines under #textbook, for instance.

        # So, everything in here is assumed to be a list. It's a little 
        # inconvenient, but I think it's worth it in case people edit the
        # markdown/txt files by hand. OTOH, just saying "put in multiple lines"
        # would work and simplify this a bit.
        self._scientistdict = {}
        self._fname = fname
        for (k,v) in scientistparts(fname):
            if k in self._scientistdict:
                self._scientistdict[k].append(v)
            else:
                self._scientistdict[k] = [v]
        self.name = self._scientistdict['Name'][0][0].strip()
    def matchestextbook(self,textbookname,verbose=False):
        # Should match title. For now, just match edition
        # If matches, return chapter and section
        # This is why I need a dropdown for "known" textbooks. And I need 
        # to fix the others via a script.
        # Right now, we only know Knight 3rd edition (that's calc)
        # and Knight 2nd edition (that's alg.)
        for entry in self._scientistdict['Textbook']:
            for txt in entry:
                if txt.startswith(textbookname):
                    parts = txt.replace(textbookname,'').split(',')
                    chapter,section = None,None
                    for part in [p for p in parts if p]:
                        _part = part.lower().strip()
                        if _part.startswith('chapter'):
                            chapter = int(_part.replace('chapter','').strip())
                        elif _part.startswith('section'):
                            section = int(_part.replace('section','').strip())
                        else:
                            print('Unknown textbook part',part)
                    return({'Chapter':chapter,
                            'Section':section})
        return False
    def todf(self,chapter,section):
        # See comments in __init__ but things may be nested one
        # level deeper than you expect.
        parts = {'Chapter':chapter,
                'Section':section}
        for (k,v) in self._scientistdict.items():
            if k == 'Textbook':
                continue
            elif k == 'Photo':
                # Can be multiple photos. 
                # Each gets a link to the original source of the photo, and we follow
                # it all up with a URL to download the slide.
                for photos in v:
                    for photo in photos:
                        this_entry = f'<a href="{photo}"><img src="{photo}" width="300"/></a>'
                        if k not in parts:
                            parts[k] = this_entry
                        else:
                            parts[k] = ' '.join([parts[k],this_entry])
            elif k == 'Sources':
                # Wrap it up in a link.
                for sources in v:
                    for source in sources:
                        formatted = ''
                        for sourcepart in source.split():
                            if sourcepart.startswith('http'):
                                linkpart = sourcepart.strip()
                                if linkpart.endswith('.'):
                                    linkpart = linkpart[:-1]
                                formatted +=  f' <a href="{linkpart}">Link</a> '
                            else:
                                formatted += ' ' + sourcepart
                        if 'Sources' in parts:
                            parts['Sources'] += formatted
                        else:
                            parts['Sources'] = formatted
            else:
                # v looks like [['Myname']] or [['Description','Description','more description']]
                for entry in v:
                    this_entry = ' '.join(entry)
                if k not in parts:
                    parts[k] = this_entry
                else:
                    parts[k] = ' '.join([parts[k],this_entry])
        if 'Photo' in parts:
            slide_link = '<a href="Textbooks/{n}.pptx">Download Slide</a>'.format(
                        n=parts['Name'],)
            parts['Photo'] = parts['Photo'] + ' ' + slide_link
        #for (k,v) in parts.items():
        #    print(f'{k}: {v}')
        df = pd.DataFrame(data=[parts])
        return df


def add_scientist_slide(prs,scientist,verbose=False):
    """Creates image file as side effect
    """
    for photos in scientist._scientistdict['Photo']:
        for photo in photos:
            img_path = 'test.jpg'
            if verbose:
                print(f'Grabbing image: {photo}')
            img_data = requests.get(photo).content
            with open(img_path, 'wb') as handler:
                handler.write(img_data)
            blank_slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(blank_slide_layout)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.text = scientist.name

            top = Inches(1.75)
            left = Inches(0.5)
            height = Inches(5)
            pic = slide.shapes.add_picture(img_path, left, top, height=height)

            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = open(scientist._fname).read()

def maketextbook(fname,*,textbookname,scientists,verbose=False):
    """This also makes the associated PPTx slides.
    """
    textbook = pd.read_csv(fname)
    prs0 = Presentation() # The whole textbook. Individual slides made separately below.
    for scientist in scientists:
        match = scientist.matchestextbook(textbookname)
        if match:
            chapter,section = match['Chapter'],match['Section']
            if verbose:
                print('Adding',scientist._scientistdict['Name'])
            df = scientist.todf(chapter=chapter,section=section)
            textbook = textbook.append(df)
            
            if 'Photo' in scientist._scientistdict:
                prs = Presentation()
                add_scientist_slide(prs0,scientist,verbose=verbose)
                add_scientist_slide(prs,scientist,verbose=verbose)
                # get the original, unparsed photos bit.
                prs.save('Textbooks/{n}.pptx'.format(n=scientist.name))


    # I want the order of columns
    cols = ['Chapter','Section','Topics','Photo','Name','Description','Sources']
    cols = cols + [i for i in textbook.columns.tolist() if i not in cols]
    textbook = textbook[cols]
    textbook = textbook.sort_values(by=['Chapter','Section','Topics','Description','Name'],
                                    ascending=[True, True, True, False, False])
    textbook = textbook.fillna('').set_index(['Chapter', 'Section','Topics'])

    # Github doesn't quite do what I expect with the table
    # rendering. Specifically, I can't seem to just specify the size of
    # a photo either in the `img` tag or in the `th` element. I've
    # tried both CSS and HTML, but what seems to actually work is
    # setting it in HTML or CSS ... and then making the column header
    # wider.
    spaces = '&nbsp;'*50
    textbook.rename(columns={'Photo':'Photo'+spaces,'Description':'Description'+spaces+spaces}, inplace=True)
    prs0.save(f'Textbooks/{textbookname}.pptx')    
    return textbook


def makeKnightCalc3rd():
    """Creates all of the PPT files as a side-effect, returns the HTML
    """
    scientists = [Scientist(fname) for fname in glob.glob('Scientists/*.txt')]
    textbook = maketextbook('Textbooks/Knight3rdEdition.csv',
                            textbookname=known_textbooks[0],
                            scientists=scientists)
    with pd.option_context('display.max_colwidth', -1):
        output_html = textbook.to_html(escape=False)

    f = open('knight.html','w')
    f.write(output_html)
    f.close()
    return(output_html)

